import os
import json
from typing import List, Tuple

import numpy as np
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import linear_kernel

try:
    from pypdf import PdfReader
except ImportError:  # pragma: no cover
    PdfReader = None

try:
    from google.cloud import vision
    from google.cloud import storage
except ImportError:  # pragma: no cover
    vision = None
    storage = None

try:
    from docx import Document
except ImportError:  # pragma: no cover
    Document = None

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover
    Presentation = None

import email
from email import policy
from email.parser import BytesParser

from app import db

CHUNK_WORDS = 400
CHUNK_OVERLAP = 50
TOP_K_DEFAULT = 8


def save_uploaded_file(file_obj, storage_dir, filename_hint):
    os.makedirs(storage_dir, exist_ok=True)
    safe_name = filename_hint.replace(" ", "_")
    path = os.path.join(storage_dir, safe_name)
    with open(path, "wb") as f:
        f.write(file_obj.read())
    return path


def extract_pdf_pages(file_path: str) -> List[Tuple[int, str]]:
    if PdfReader is None:
        raise RuntimeError("pypdf is not installed")

    reader = PdfReader(file_path)
    pages = []
    for idx, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        text = " ".join(text.split())
        if text:
            pages.append((idx, text))
    return pages


def extract_pdf_pages_with_ocr(file_path: str) -> List[Tuple[int, str]]:
    if vision is None or storage is None:
        raise RuntimeError("google-cloud-vision/storage not installed")
    bucket_name = os.getenv("GCP_OCR_BUCKET")
    if not bucket_name:
        raise RuntimeError("GCP_OCR_BUCKET is not set")

    vision_client = vision.ImageAnnotatorClient()
    storage_client = storage.Client()
    bucket = storage_client.bucket(bucket_name)

    filename = os.path.basename(file_path)
    upload_blob = bucket.blob(f"temp-uploads/{os.urandom(8).hex()}-{filename}")
    upload_blob.upload_from_filename(file_path)

    gcs_source_uri = f"gs://{bucket_name}/{upload_blob.name}"
    gcs_dest_prefix = f"ocr-results/{os.urandom(8).hex()}-"
    gcs_dest_uri = f"gs://{bucket_name}/{gcs_dest_prefix}"

    request = {
        "requests": [
            {
                "input_config": {"gcs_source": {"uri": gcs_source_uri}, "mime_type": "application/pdf"},
                "features": [{"type_": vision.Feature.Type.DOCUMENT_TEXT_DETECTION}],
                "output_config": {"gcs_destination": {"uri": gcs_dest_uri}, "batch_size": 100},
            }
        ]
    }

    operation = vision_client.async_batch_annotate_files(request)
    operation.result(timeout=600)

    pages = []
    for blob in storage_client.list_blobs(bucket_name, prefix=gcs_dest_prefix):
        content = blob.download_as_bytes()
        result = json.loads(content.decode("utf-8"))
        responses = result.get("responses", [])
        for idx, response in enumerate(responses, start=1):
            full = response.get("fullTextAnnotation", {}).get("text", "")
            full = " ".join(full.split())
            if full:
                pages.append((idx, full))
        blob.delete()

    upload_blob.delete()
    return pages


def extract_txt_pages(file_path: str) -> List[Tuple[int, str]]:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    text = " ".join(text.split())
    return [(1, text)] if text else []


def extract_eml_pages(file_path: str) -> List[Tuple[int, str]]:
    with open(file_path, "rb") as f:
        msg = BytesParser(policy=policy.default).parse(f)

    parts = []
    subject = msg.get("subject", "")
    sender = msg.get("from", "")
    date = msg.get("date", "")
    if subject:
        parts.append(f"Subject: {subject}")
    if sender:
        parts.append(f"From: {sender}")
    if date:
        parts.append(f"Date: {date}")

    body_texts = []
    if msg.is_multipart():
        for part in msg.walk():
            if part.get_content_type() == "text/plain":
                try:
                    body_texts.append(part.get_content())
                except Exception:
                    pass
    else:
        try:
            body_texts.append(msg.get_content())
        except Exception:
            pass

    if body_texts:
        parts.append("\n".join(body_texts))

    text = " ".join(" ".join(parts).split())
    return [(1, text)] if text else []


def extract_docx_pages(file_path: str) -> List[Tuple[int, str]]:
    if Document is None:
        raise RuntimeError("python-docx is not installed")
    doc = Document(file_path)
    text = " ".join(p.text for p in doc.paragraphs if p.text)
    text = " ".join(text.split())
    return [(1, text)] if text else []


def extract_pptx_pages(file_path: str) -> List[Tuple[int, str]]:
    if Presentation is None:
        raise RuntimeError("python-pptx is not installed")
    pres = Presentation(file_path)
    pages = []
    for idx, slide in enumerate(pres.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        combined = " ".join(" ".join(texts).split())
        if combined:
            pages.append((idx, combined))
    return pages


def detect_doc_type(filename: str) -> str:
    lower = filename.lower()
    if lower.endswith(".pdf"):
        return "pdf"
    if lower.endswith(".txt"):
        return "txt"
    if lower.endswith(".eml"):
        return "eml"
    if lower.endswith(".docx"):
        return "docx"
    if lower.endswith(".pptx"):
        return "pptx"
    return "unknown"


def _chunk_text(text: str, words_per_chunk: int = CHUNK_WORDS, overlap: int = CHUNK_OVERLAP) -> List[str]:
    words = text.split()
    if not words:
        return []
    chunks = []
    start = 0
    while start < len(words):
        end = min(len(words), start + words_per_chunk)
        chunk = " ".join(words[start:end])
        chunks.append(chunk)
        if end == len(words):
            break
        start = max(0, end - overlap)
    return chunks


def build_chunks(pages: List[Tuple[int, str]]):
    chunks = []
    for page_num, text in pages:
        for idx, chunk in enumerate(_chunk_text(text)):
            chunks.append(
                {
                    "page_num": page_num,
                    "chunk_index": idx,
                    "chunk_text": chunk,
                }
            )
    return chunks


def store_document(conn, file_obj, filename, doc_type, source_type, is_global=False, company_id=None, actor="user", use_ocr=False):
    storage_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "data", "uploads")
    path = save_uploaded_file(file_obj, storage_dir, filename)

    if doc_type == "pdf":
        pages = extract_pdf_pages(path)
        if not pages and use_ocr:
            pages = extract_pdf_pages_with_ocr(path)
    elif doc_type == "txt":
        pages = extract_txt_pages(path)
    elif doc_type == "eml":
        pages = extract_eml_pages(path)
    elif doc_type == "docx":
        pages = extract_docx_pages(path)
    elif doc_type == "pptx":
        pages = extract_pptx_pages(path)
    else:
        pages = []

    full_text = " ".join(text for _, text in pages)
    doc_id = db.create_document(
        conn,
        filename,
        path,
        is_global=is_global,
        actor=actor,
        doc_type=doc_type,
        source_type=source_type,
        extracted_text=full_text,
    )
    if company_id is not None:
        db.link_document_to_company(conn, doc_id, company_id, actor=actor)

    chunks = build_chunks(pages)
    db.add_document_chunks(conn, doc_id, chunks)
    return doc_id, len(pages), len(chunks)


def store_pdf(conn, file_obj, filename, is_global=False, company_id=None, actor="user"):
    return store_document(conn, file_obj, filename, "pdf", "pdf", is_global=is_global, company_id=company_id, actor=actor)


def build_query_from_signals(signal_defs):
    parts = []
    for d in signal_defs:
        if d["disabled"]:
            continue
        name = d["name"] or ""
        desc = d["description"] or ""
        parts.append(f"{name} {desc}".strip())
    base = " ".join(parts)
    guidance = "startup product traction founders team funding stage market pricing"
    return f"{base} {guidance}".strip()


def retrieve_top_chunks(chunks, query, top_k=TOP_K_DEFAULT):
    if not chunks:
        return []
    texts = [c["chunk_text"] for c in chunks]
    vectorizer = TfidfVectorizer(stop_words="english", max_features=5000)
    tfidf = vectorizer.fit_transform(texts)
    query_vec = vectorizer.transform([query])
    scores = linear_kernel(query_vec, tfidf).flatten()
    top_idx = np.argsort(scores)[::-1][:top_k]
    return [chunks[i] for i in top_idx if scores[i] > 0]


def build_context(chunks):
    sections = []
    for ch in chunks:
        filename = ch.get("filename") if isinstance(ch, dict) else None
        if filename:
            sections.append(f"[Page {ch['page_num']}] ({filename}) {ch['chunk_text']}")
        else:
            sections.append(f"[Page {ch['page_num']}] {ch['chunk_text']}")
    return "\n\n".join(sections)


def get_company_chunks(conn, company_id, source_types=None):
    rows = db.list_document_chunks_for_company(conn, company_id, source_types=source_types)
    return [dict(r) for r in rows]
